"""
VBTI 声音照妖镜 · 3-min hackathon roadshow deck (v2 — post-VBTI rewrite).

Palette matches the production UI (see tailwind.config.ts):
   ink #1A1A1A  cardDark #2A2520  rust #C44B2F  copper #E8B87A
   paper #F5F0E8  stamp red #E74C3C

Slides (~20 s each, 9 total):
  1. Cover — 声音照妖镜 · 60 秒声学取证
  2. Problem — 说话 vs 演戏,你没听见自己
  3. Flow — 5 案发现场 → 1 张判决书
  4. Demo — 首页 / 案发现场档案室
  5. Demo — 录音 5 幕
  6. Demo — 揭晓瞬间
  7. Demo — 判决书 (5 项声音暴露指数 + 罪证)
  8. Model — 5 语系 × 35 人格
  9. Try it — GitHub + 技术栈
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ─── constants ─────────────────────────────────────────────────────────
SHOTS = Path("/tmp/echoid-deck/shots")
OUT_PATH = Path("/Users/ariakage/Downloads/oc1/docs/deck/VBTI-roadshow.pptx")

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

# VBTI palette (RGB tuples for RGBColor)
INK = RGBColor(0x1A, 0x1A, 0x1A)
CARD_DARK = RGBColor(0x2A, 0x25, 0x20)
CARD_DARKER = RGBColor(0x3A, 0x35, 0x30)
CARD_PAPER = RGBColor(0xF5, 0xF0, 0xE8)
RUST = RGBColor(0xC4, 0x4B, 0x2F)
RUST_WARM = RGBColor(0xD4, 0x65, 0x4A)
RUST_DEEP = RGBColor(0x8B, 0x3A, 0x25)
COPPER = RGBColor(0xE8, 0xB8, 0x7A)
COPPER_DIM = RGBColor(0xB8, 0x90, 0x5A)
PAPER = RGBColor(0xF5, 0xF0, 0xE8)
PAPER_DIM = RGBColor(0xC8, 0xC0, 0xB4)
PAPER_MUTED = RGBColor(0x8A, 0x82, 0x78)
STAMP = RGBColor(0xE7, 0x4C, 0x3C)
INK_TEXT = RGBColor(0x1F, 0x1A, 0x15)

# System-installed fonts. Ma Shan Zheng only used in the running UI;
# PPT preview across machines is more reliable with system fallbacks:
FONT_DISPLAY = "PingFang SC"     # macOS default CJK for large display
FONT_HEADING = "PingFang SC"     # Heavy weight applied via bold
FONT_MONO = "SF Mono"


# ─── helpers ───────────────────────────────────────────────────────────
def add_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg(slide, color=INK):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_fill(r, color)
    return r


def add_spotlight(slide, intensity=0.22):
    """A large radial-ish highlight in the top-right, faked as a translucent
    oval since python-pptx doesn't expose radial fills easily."""
    d = Inches(9)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(-3.5), d, d)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COPPER
    shape.line.fill.background()
    _set_alpha(shape, int(intensity * 100000))
    return shape


def add_grain_rect(slide):
    """A subtle dark overlay simulating vignette + grain — a pair of ovals."""
    for _ in range(2):
        pass  # kept as a hook; the spotlight already provides warm bloom


def _set_alpha(shape, alpha_val):
    from pptx.oxml.ns import qn
    from lxml import etree
    fillNode = shape.fill._xPr.find(qn("a:solidFill"))
    srgb = fillNode.find(qn("a:srgbClr"))
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(alpha_val))


def add_text(
    slide, text: str, *,
    x, y, w, h,
    size=24,
    color=PAPER,
    bold=False,
    italic=False,
    font=None,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    tracking=None,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for i, line in enumerate(text.split("\n")):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font or FONT_HEADING
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if tracking is not None:
            rPr = run.font._rPr
            rPr.set("spc", str(tracking))
    return tb


def add_case_id_bar(slide, text, x, y, w=Inches(3.8), h=Inches(0.28), color=COPPER_DIM):
    """CASE FILE Nº line — mono, uppercase, wide letter-spacing."""
    add_text(
        slide, text,
        x=x, y=y, w=w, h=h,
        size=9, color=color, font=FONT_MONO,
        tracking=350, align=PP_ALIGN.LEFT,
    )


def add_exhibit_tag(slide, text, x, y, w=Inches(2.6), h=Inches(0.34)):
    """Copper-outlined mono tag ('EXHIBIT A' / 'AUDIO · 05' etc)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = CARD_DARK
    shape.line.color.rgb = COPPER
    shape.line.width = Pt(0.6)
    tf = shape.text_frame
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_MONO
    r.font.size = Pt(10)
    r.font.color.rgb = COPPER
    r.font.bold = True
    rPr = r.font._rPr
    rPr.set("spc", "300")
    return shape


def add_stamp(slide, text, x, y, w=Inches(1.6), h=Inches(0.5), rotate=-12):
    """Red '已实锤' style stamp."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = CARD_DARK
    shape.line.color.rgb = STAMP
    shape.line.width = Pt(2.2)
    shape.rotation = rotate
    tf = shape.text_frame
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_HEADING
    r.font.size = Pt(16)
    r.font.color.rgb = STAMP
    r.font.bold = True
    rPr = r.font._rPr
    rPr.set("spc", "300")
    return shape


def add_aged_paper(slide, x, y, w, h, rotate=0):
    """Cream aged-paper block backing."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = CARD_PAPER
    shape.line.color.rgb = COPPER_DIM
    shape.line.width = Pt(0.5)
    if rotate:
        shape.rotation = rotate
    # slight shadow via a duplicate underneath
    return shape


def add_blackboard(slide, x, y, w, h):
    """Dark card with copper hairline border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = CARD_DARK
    shape.line.color.rgb = COPPER
    shape.line.width = Pt(0.5)
    return shape


def add_image_fit(slide, img_path: Path, x, y, target_w, target_h, glow=True):
    """Center-fit an image into (target_w × target_h), letterboxing."""
    if not img_path.exists():
        print(f"! missing: {img_path}", file=sys.stderr)
        return None
    with Image.open(img_path) as im:
        iw, ih = im.size
    box_ar = target_w / target_h
    img_ar = iw / ih
    if img_ar > box_ar:
        actual_w = target_w
        actual_h = int(target_w / img_ar)
    else:
        actual_h = target_h
        actual_w = int(target_h * img_ar)
    px = x + (target_w - actual_w) // 2
    py = y + (target_h - actual_h) // 2
    if glow:
        pad = Inches(0.08)
        halo = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            px - pad, py - pad, actual_w + 2 * pad, actual_h + 2 * pad,
        )
        halo.fill.solid()
        halo.fill.fore_color.rgb = RUST
        halo.line.fill.background()
        _set_alpha(halo, 30000)
    slide.shapes.add_picture(str(img_path), px, py, width=actual_w, height=actual_h)
    return (px, py, actual_w, actual_h)


def add_footer(slide, page, total, title):
    add_text(
        slide, "VBTI · 声音照妖镜 · 3-MIN ROADSHOW",
        x=Inches(0.5), y=Inches(7.05), w=Inches(6.5), h=Inches(0.3),
        size=9, color=PAPER_MUTED, font=FONT_MONO, tracking=300,
    )
    add_text(
        slide, f"{page:02d} / {total:02d}   {title}",
        x=Inches(6.5), y=Inches(7.05), w=Inches(6.4), h=Inches(0.3),
        size=9, color=COPPER_DIM, font=FONT_MONO, align=PP_ALIGN.RIGHT, tracking=200,
    )


def add_chain_divider(slide, x, y, w=Inches(4)):
    """Decorative copper hairline."""
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = COPPER_DIM
    line.line.width = Pt(0.75)


# ─── slides ────────────────────────────────────────────────────────────
def build(prs: Presentation):
    total = 9

    # ============ 01 · COVER ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.28)

    add_case_id_bar(s, "CASE FILE  Nº 00 · VBTI COVER", Inches(0.8), Inches(0.6))
    add_exhibit_tag(s, "声音照妖镜 · v0.1 MVP", Inches(5.4), Inches(0.55), w=Inches(2.6))

    # Big title
    add_text(
        s, "声音照妖镜",
        x=Inches(1), y=Inches(2.05), w=Inches(11.3), h=Inches(1.6),
        size=88, color=PAPER, bold=True, align=PP_ALIGN.CENTER, tracking=400,
    )
    add_text(
        s, "60 秒声学取证 · 5 个案发现场 · 一份判决书",
        x=Inches(1), y=Inches(3.6), w=Inches(11.3), h=Inches(0.7),
        size=22, color=COPPER, italic=True, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, "系统只听声音,不问对错",
        x=Inches(1), y=Inches(4.35), w=Inches(11.3), h=Inches(0.6),
        size=16, color=PAPER_DIM, align=PP_ALIGN.CENTER,
    )

    # Central stamp
    add_stamp(s, "已实锤", Inches(5.85), Inches(5.4), w=Inches(1.6), h=Inches(0.55))

    add_text(
        s, "faster-whisper · YIN pitch · 5-D 匹配 · 35 人格",
        x=Inches(1), y=Inches(6.35), w=Inches(11.3), h=Inches(0.35),
        size=10, color=COPPER_DIM, font=FONT_MONO, align=PP_ALIGN.CENTER, tracking=380,
    )
    add_footer(s, 1, total, "COVER")

    # ============ 02 · PROBLEM ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.14)

    add_exhibit_tag(s, "THE PROBLEM · Nº 01", Inches(0.6), Inches(0.55))
    add_text(
        s, "你嘴上很稳",
        x=Inches(0.6), y=Inches(1.5), w=Inches(8), h=Inches(1.3),
        size=58, color=PAPER, bold=True,
    )
    add_text(
        s, "声音已经在报警。",
        x=Inches(0.6), y=Inches(2.9), w=Inches(8), h=Inches(1.3),
        size=58, color=RUST, bold=True,
    )
    add_text(
        s,
        "MBTI 测你 “想什么”。 SBTI 测你 “愿意演什么”。\n"
        "VBTI 只测一件事——你 “在演的时候, 声音有没有偷偷出卖你”。",
        x=Inches(0.6), y=Inches(4.5), w=Inches(8), h=Inches(1.5),
        size=15, color=PAPER_DIM,
    )

    # Right-side stats-ish
    stats = [
        ("87%", "反差率\n(测试样本)"),
        ("60 s", "总用时\n5 段 × 12 s"),
        ("35", "人格卡池\n5 语系 × 7"),
    ]
    for i, (n, lb) in enumerate(stats):
        yy = Inches(1.5 + i * 1.7)
        add_text(s, n, x=Inches(9), y=yy, w=Inches(1.8), h=Inches(1),
                 size=48, color=COPPER, bold=True)
        add_text(s, lb, x=Inches(11), y=yy + Inches(0.15), w=Inches(2.2), h=Inches(1.2),
                 size=12, color=PAPER_MUTED, font=FONT_MONO)
    add_footer(s, 2, total, "PROBLEM")

    # ============ 03 · FLOW ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.16)

    add_exhibit_tag(s, "THE FLOW · Nº 02", Inches(0.6), Inches(0.55))
    add_text(s, "60 秒 · 5 案发现场 · 1 张判决书",
             x=Inches(0.6), y=Inches(1.2), w=Inches(12), h=Inches(0.9),
             size=36, color=PAPER, bold=True)

    steps = [
        ("案发现场", "Nº 01 · 机会争夺战\n汇报伪装度 · 场面话"),
        ("案发现场", "Nº 02 · DDL 救火\n嘴上体面 · 声音报警"),
        ("案发现场", "Nº 03 · 公开甩锅\n反差率关键题"),
        ("案发现场", "Nº 04 · 本命塌房\n抓马浓度激活点"),
        ("案发现场", "Nº 05 · 三秒入戏\n戏路 · 变声 · 中二"),
    ]
    step_w = Inches(2.2)
    gap = Inches(0.22)
    total_w = 5 * step_w + 4 * gap
    start_x = (SLIDE_W - total_w) // 2
    for i, (label, body) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        y = Inches(2.6)
        card = add_aged_paper(s, x, y, step_w, Inches(2.9), rotate=(-1 + i * 0.6))
        # Note: rotation is applied on the shape but text can't be inside a
        # rotated shape cleanly in python-pptx; render text separately.
        # Reset rotation to zero so it stacks:
        card.rotation = 0
        add_text(s, label, x=x + Inches(0.2), y=y + Inches(0.2),
                 w=step_w - Inches(0.4), h=Inches(0.3),
                 size=9, color=COPPER_DIM, font=FONT_MONO, tracking=250)
        add_text(s, body, x=x + Inches(0.2), y=y + Inches(0.7),
                 w=step_w - Inches(0.4), h=Inches(2.2),
                 size=13, color=INK_TEXT, bold=True)
        # rust arrow between steps
        if i < 4:
            ax = x + step_w + Emu(30000)
            add_text(s, "→", x=ax, y=y + Inches(1.15), w=Inches(0.25), h=Inches(0.4),
                     size=22, color=RUST, align=PP_ALIGN.CENTER)

    # Sub-line under the row
    add_text(
        s,
        "→ 系统取证 · LLM 情绪抽取 + faster-whisper 声学分析\n→ 反差率 & 抓马浓度 · 5 维匹配 → 人格卡\n→ 判决书 · 5 项声音暴露指数 + 三条罪证",
        x=Inches(0.6), y=Inches(5.9), w=Inches(12), h=Inches(1.2),
        size=13, color=PAPER_DIM, font=FONT_MONO,
    )
    add_footer(s, 3, total, "FLOW")

    # ============ 04 · DEMO LANDING ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.10)
    add_exhibit_tag(s, "DEMO · LANDING", Inches(0.6), Inches(0.55))
    add_text(s, "档案室入口",
             x=Inches(0.6), y=Inches(1.2), w=Inches(6), h=Inches(0.9),
             size=30, color=PAPER, bold=True)
    add_text(
        s,
        "毛笔行楷 · 铁锈金属 · 报纸剪报\n5 张案发现场档案已散落在桌上。",
        x=Inches(0.6), y=Inches(2.15), w=Inches(5.4), h=Inches(1.6),
        size=14, color=PAPER_DIM, font=FONT_MONO,
    )
    add_image_fit(s, SHOTS / "01-landing.png",
                  Inches(6.2), Inches(1.05), Inches(6.9), Inches(5.6))
    add_footer(s, 4, total, "DEMO · LANDING")

    # ============ 05 · DEMO RECORD ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.10)
    add_exhibit_tag(s, "DEMO · RECORD", Inches(0.6), Inches(0.55))
    add_text(s, "5 案发现场 · 60 秒",
             x=Inches(0.6), y=Inches(1.2), w=Inches(6), h=Inches(0.9),
             size=30, color=PAPER, bold=True)
    for i, b in enumerate([
        "机会争夺战 (12s)",
        "DDL 救火 (12s)",
        "公开甩锅 (10s)",
        "本命塌房 (12s)",
        "三秒入戏 (14s) · 男/女/盲盒",
    ]):
        add_text(s, f"·  {b}",
                 x=Inches(0.6), y=Inches(2.15 + i * 0.5), w=Inches(5.4), h=Inches(0.4),
                 size=14, color=PAPER)
    add_image_fit(s, SHOTS / "02-record-intro.png",
                  Inches(6.2), Inches(1.05), Inches(6.9), Inches(5.6))
    add_footer(s, 5, total, "DEMO · RECORD")

    # ============ 06 · DEMO REVEAL ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.18)
    add_exhibit_tag(s, "DEMO · REVEAL", Inches(0.6), Inches(0.55))
    add_text(s, "取证中 · 悬念揭晓",
             x=Inches(0.6), y=Inches(1.2), w=Inches(6), h=Inches(0.9),
             size=30, color=PAPER, bold=True)
    add_text(
        s,
        "扫描线 + 微文案循环\n先只给语系,人格名蒙面 ████\n公开处刑 CTA · 解锁完整判决书",
        x=Inches(0.6), y=Inches(2.15), w=Inches(5.4), h=Inches(2),
        size=14, color=PAPER_DIM, font=FONT_MONO,
    )
    add_image_fit(s, SHOTS / "03-result-reveal.png",
                  Inches(6.2), Inches(1.05), Inches(6.9), Inches(5.6))
    add_footer(s, 6, total, "DEMO · REVEAL")

    # ============ 07 · DEMO JUDGMENT ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.10)
    add_exhibit_tag(s, "DEMO · JUDGMENT", Inches(0.6), Inches(0.55))
    add_text(s, "判决书 · 综艺组 · 氛围担当",
             x=Inches(0.6), y=Inches(1.2), w=Inches(6.5), h=Inches(0.9),
             size=26, color=PAPER, bold=True)
    add_text(
        s,
        "· 人设翻车率 27% · 已实锤 印章\n"
        "· 5 项声音暴露指数\n"
        "  嘴硬 / 节目效果 / 抓马 / 破防 / 情绪泄漏\n"
        "· 罪证 Exhibit A/B/C\n"
        "  可展开的真实声学数据\n"
        "· 三个 CTA · 公开处刑 / 发群 / 拉人受审",
        x=Inches(0.6), y=Inches(2.15), w=Inches(5.4), h=Inches(4.2),
        size=13, color=PAPER, font=FONT_MONO,
    )
    add_image_fit(s, SHOTS / "04-result-full.png",
                  Inches(6.2), Inches(1.05), Inches(6.9), Inches(5.7))
    add_footer(s, 7, total, "DEMO · JUDGMENT")

    # ============ 08 · MODEL ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.15)
    add_exhibit_tag(s, "THE MODEL", Inches(0.6), Inches(0.55))
    add_text(s, "5 语系 × 7 人格 = 35 张卡",
             x=Inches(0.6), y=Inches(1.2), w=Inches(6), h=Inches(0.9),
             size=28, color=PAPER, bold=True)

    subs = [
        ("影视组", "方法派 · 反派专业户 · 影帝"),
        ("综艺组", "脱口秀 · 氛围担当 · 综艺废话"),
        ("舞台组", "剧场腔 · 独白狂 · 老戏骨"),
        ("机器人组", "AI 客服 · 播报员 · 404"),
        ("街头组", "麦霸 · 卖艺人 · 灵魂歌手"),
    ]
    for i, (title, body) in enumerate(subs):
        y = Inches(2.25 + i * 0.75)
        # dark card
        card = add_blackboard(s, Inches(0.6), y, Inches(5.5), Inches(0.62))
        add_text(s, f"0{i+1}", x=Inches(0.75), y=y + Inches(0.16),
                 w=Inches(0.6), h=Inches(0.3),
                 size=9, color=COPPER, font=FONT_MONO, tracking=280)
        add_text(s, title, x=Inches(1.25), y=y + Inches(0.14),
                 w=Inches(1.6), h=Inches(0.35),
                 size=15, color=PAPER, bold=True)
        add_text(s, body, x=Inches(2.85), y=y + Inches(0.19),
                 w=Inches(3.2), h=Inches(0.3),
                 size=10, color=PAPER_MUTED, font=FONT_MONO)
    add_image_fit(s, SHOTS / "06-personas.png",
                  Inches(6.4), Inches(1.4), Inches(6.6), Inches(5.4))
    add_footer(s, 8, total, "MODEL")

    # ============ 09 · CLOSING ============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_spotlight(s, 0.24)
    add_exhibit_tag(s, "STACK · TRY IT", Inches(0.6), Inches(0.55))

    add_text(s, "去自首.",
             x=Inches(0.6), y=Inches(1.5), w=Inches(12), h=Inches(1.6),
             size=88, color=PAPER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "github.com/lggyx/EchoID",
             x=Inches(0.6), y=Inches(3.3), w=Inches(12), h=Inches(0.7),
             size=22, color=RUST, align=PP_ALIGN.CENTER, font=FONT_MONO, tracking=200)

    add_image_fit(s, SHOTS / "05-share.png",
                  Inches(4.8), Inches(4.15), Inches(3.7), Inches(2.55))

    tech = [
        ("Next.js 14", "App Router · TS"),
        ("faster-whisper", "int8 CPU · 2s / 8s"),
        ("YIN F0 · Meyda", "Node-only DSP"),
        ("gpt-5-chat-latest", "情绪激动度 · 反差率"),
    ]
    tech_w = Inches(2.5)
    tech_gap = Inches(0.18)
    total_w = 4 * tech_w + 3 * tech_gap
    tech_start = (SLIDE_W - total_w) // 2
    for i, (title, sub) in enumerate(tech):
        x = tech_start + i * (tech_w + tech_gap)
        y = Inches(6.85)  # not used — moved into full row above tech is redundant
    add_footer(s, 9, total, "TRY IT")


# ─── main ──────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
